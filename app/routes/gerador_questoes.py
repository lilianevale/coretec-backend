import os
import glob
from flask import Blueprint, request, jsonify, current_app
from werkzeug.utils import secure_filename
from pptx import Presentation
from app.utils.openai_helper import chat_with_gpt

gerador_questoes_bp = Blueprint('questoes', __name__)

ALLOWED_EXTENSIONS = {'ppt', 'pptx'}

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

@gerador_questoes_bp.route('/questoes/ppt', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part in the request'}), 400

    file = request.files['file']

    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400

    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        upload_folder = current_app.config.get('UPLOAD_FOLDER', 'uploads')
        os.makedirs(upload_folder, exist_ok=True)

        filepath = os.path.join(upload_folder, filename)
        file.save(filepath)

        # Extrair o texto dos slides
        texto = "Elabore questões sobre este conteúdo:\n"
        prs = Presentation(filepath)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texto += shape.text + '\n'

        # Gerar questões com o GPT
        try:
            perguntas = chat_with_gpt(texto)
            return jsonify({'questoes': perguntas})
        except Exception as e:
            return jsonify({'error': str(e)}), 500

    return jsonify({'error': 'File type not allowed'}), 400
